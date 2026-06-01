# ==================== ШАБЛОНЫ ДЛЯ ПРЕМИУМ (с картинками) ====================

TEMPLATE1 = """
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ── Титульный слайд (тёмный + акцент слева) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(13, 17, 30)

left = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.5), Inches(7.5))
left.fill.solid(); left.fill.fore_color.rgb = RGBColor(99, 102, 241); left.line.fill.background()

bottom = slide.shapes.add_shape(1, Inches(0), Inches(6.8), Inches(10), Inches(0.7))
bottom.fill.solid(); bottom.fill.fore_color.rgb = RGBColor(30, 27, 75); bottom.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(8.5), Inches(2.0))
tf = title_box.text_frame; tf.word_wrap = True
tf.text = "НАЗВАНИЕ ПРЕЗЕНТАЦИИ"
p = tf.paragraphs[0]; p.font.size = Pt(46); p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

sep = slide.shapes.add_shape(1, Inches(0.9), Inches(3.8), Inches(3), Inches(0.07))
sep.fill.solid(); sep.fill.fore_color.rgb = RGBColor(99,102,241); sep.line.fill.background()

sub_box = slide.shapes.add_textbox(Inches(0.9), Inches(4.0), Inches(8), Inches(0.8))
sf = sub_box.text_frame
sf.text = "Подзаголовок"
sf.paragraphs[0].font.size = Pt(20); sf.paragraphs[0].font.color.rgb = RGBColor(165, 180, 252)

# ── Контентный слайд (текст слева, картинка справа) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(13, 17, 30)

left = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.5), Inches(7.5))
left.fill.solid(); left.fill.fore_color.rgb = RGBColor(99, 102, 241); left.line.fill.background()

header = slide.shapes.add_shape(1, Inches(0.5), Inches(0), Inches(9.5), Inches(1.1))
header.fill.solid(); header.fill.fore_color.rgb = RGBColor(23, 25, 45); header.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.12), Inches(8), Inches(0.85))
tf = title_box.text_frame
tf.text = "ЗАГОЛОВОК СЛАЙДА"
p = tf.paragraphs[0]; p.font.size = Pt(28); p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

text_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.3), Inches(4.6), Inches(5.9))
tf2 = text_box.text_frame; tf2.word_wrap = True
tf2.text = "Первый пункт"
tf2.paragraphs[0].font.size = Pt(16); tf2.paragraphs[0].font.color.rgb = RGBColor(226, 232, 240)
for txt in ["Второй пункт", "Третий пункт"]:
    p = tf2.add_paragraph(); p.text = txt
    p.font.size = Pt(16); p.font.color.rgb = RGBColor(226, 232, 240); p.space_before = Pt(12)

if os.path.exists("img_0.jpg"):
    slide.shapes.add_picture("img_0.jpg", Inches(5.7), Inches(1.2), width=Inches(4.1), height=Inches(6.0))

prs.save("presentation.pptx")
"""

TEMPLATE2 = """
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ── Титульный слайд (светлый, картинка на фоне слева) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(248, 250, 252)

left_bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(4.8), Inches(7.5))
left_bg.fill.solid(); left_bg.fill.fore_color.rgb = RGBColor(15, 23, 42); left_bg.line.fill.background()

if os.path.exists("img_0.jpg"):
    slide.shapes.add_picture("img_0.jpg", Inches(0), Inches(0), width=Inches(4.8), height=Inches(7.5))

right_bg = slide.shapes.add_shape(1, Inches(4.8), Inches(0), Inches(5.2), Inches(7.5))
right_bg.fill.solid(); right_bg.fill.fore_color.rgb = RGBColor(248, 250, 252); right_bg.line.fill.background()

accent = slide.shapes.add_shape(1, Inches(4.8), Inches(0), Inches(0.12), Inches(7.5))
accent.fill.solid(); accent.fill.fore_color.rgb = RGBColor(14, 165, 233); accent.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.5), Inches(2.2))
tf = title_box.text_frame; tf.word_wrap = True
tf.text = "НАЗВАНИЕ ПРЕЗЕНТАЦИИ"
p = tf.paragraphs[0]; p.font.size = Pt(38); p.font.bold = True
p.font.color.rgb = RGBColor(15, 23, 42)

sep = slide.shapes.add_shape(1, Inches(5.2), Inches(4.1), Inches(2.5), Inches(0.07))
sep.fill.solid(); sep.fill.fore_color.rgb = RGBColor(14,165,233); sep.line.fill.background()

sub_box = slide.shapes.add_textbox(Inches(5.2), Inches(4.3), Inches(4.5), Inches(0.8))
sf = sub_box.text_frame; sf.text = "Подзаголовок"
sf.paragraphs[0].font.size = Pt(18); sf.paragraphs[0].font.color.rgb = RGBColor(100, 116, 139)

# ── Контентный слайд (картинка сверху, текст снизу) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(248, 250, 252)

top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.1))
top_bar.fill.solid(); top_bar.fill.fore_color.rgb = RGBColor(14, 165, 233); top_bar.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(9), Inches(0.85))
tf = title_box.text_frame; tf.text = "ЗАГОЛОВОК СЛАЙДА"
p = tf.paragraphs[0]; p.font.size = Pt(28); p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

if os.path.exists("img_1.jpg"):
    slide.shapes.add_picture("img_1.jpg", Inches(2.5), Inches(1.2), width=Inches(5), height=Inches(3.2))

text_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(2.7))
tf2 = text_box.text_frame; tf2.word_wrap = True
tf2.text = "Первый пункт"
tf2.paragraphs[0].font.size = Pt(16); tf2.paragraphs[0].font.color.rgb = RGBColor(30, 41, 59)
for txt in ["Второй пункт", "Третий пункт"]:
    p = tf2.add_paragraph(); p.text = txt
    p.font.size = Pt(16); p.font.color.rgb = RGBColor(30, 41, 59); p.space_before = Pt(10)

prs.save("presentation.pptx")
"""

TEMPLATE3 = """
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ── Титульный слайд (зелёный насыщенный) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(5, 46, 22)

deco_top = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.22))
deco_top.fill.solid(); deco_top.fill.fore_color.rgb = RGBColor(34, 197, 94); deco_top.line.fill.background()
deco_bot = slide.shapes.add_shape(1, Inches(0), Inches(7.28), Inches(10), Inches(0.22))
deco_bot.fill.solid(); deco_bot.fill.fore_color.rgb = RGBColor(34, 197, 94); deco_bot.line.fill.background()

if os.path.exists("img_0.jpg"):
    slide.shapes.add_picture("img_0.jpg", Inches(5.8), Inches(0.22), width=Inches(4.2), height=Inches(7.06))

fade = slide.shapes.add_shape(1, Inches(4.5), Inches(0.22), Inches(1.5), Inches(7.06))
fade.fill.solid(); fade.fill.fore_color.rgb = RGBColor(5, 46, 22); fade.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.0), Inches(2.2))
tf = title_box.text_frame; tf.word_wrap = True
tf.text = "НАЗВАНИЕ ПРЕЗЕНТАЦИИ"
p = tf.paragraphs[0]; p.font.size = Pt(42); p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

sep = slide.shapes.add_shape(1, Inches(0.6), Inches(4.1), Inches(2.5), Inches(0.07))
sep.fill.solid(); sep.fill.fore_color.rgb = RGBColor(34,197,94); sep.line.fill.background()

sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(4.3), Inches(5.0), Inches(0.8))
sf = sub_box.text_frame; sf.text = "Подзаголовок"
sf.paragraphs[0].font.size = Pt(19); sf.paragraphs[0].font.color.rgb = RGBColor(134, 239, 172)

# ── Контентный слайд ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(240, 253, 244)

header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.15))
header.fill.solid(); header.fill.fore_color.rgb = RGBColor(21, 128, 61); header.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.13), Inches(9), Inches(0.88))
tf = title_box.text_frame; tf.text = "ЗАГОЛОВОК СЛАЙДА"
p = tf.paragraphs[0]; p.font.size = Pt(28); p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.0), Inches(5.9))
tf2 = text_box.text_frame; tf2.word_wrap = True
tf2.text = "Первый пункт"
tf2.paragraphs[0].font.size = Pt(16); tf2.paragraphs[0].font.color.rgb = RGBColor(20, 83, 45)
for txt in ["Второй пункт", "Третий пункт"]:
    p = tf2.add_paragraph(); p.text = txt
    p.font.size = Pt(16); p.font.color.rgb = RGBColor(20, 83, 45); p.space_before = Pt(12)

if os.path.exists("img_1.jpg"):
    slide.shapes.add_picture("img_1.jpg", Inches(5.7), Inches(1.2), width=Inches(4.1), height=Inches(6.0))

prs.save("presentation.pptx")
"""

TEMPLATE4 = """
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ── Титульный слайд (оранжево-чёрный энергичный) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(12, 10, 9)

top_half = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(3.6))
top_half.fill.solid(); top_half.fill.fore_color.rgb = RGBColor(234, 88, 12); top_half.line.fill.background()

diag = slide.shapes.add_shape(1, Inches(0), Inches(3.2), Inches(10), Inches(0.8))
diag.fill.solid(); diag.fill.fore_color.rgb = RGBColor(194, 65, 12); diag.line.fill.background()

if os.path.exists("img_0.jpg"):
    slide.shapes.add_picture("img_0.jpg", Inches(5.5), Inches(0.0), width=Inches(4.5), height=Inches(3.6))

fade = slide.shapes.add_shape(1, Inches(4.8), Inches(0), Inches(1.0), Inches(3.6))
fade.fill.solid(); fade.fill.fore_color.rgb = RGBColor(234,88,12); fade.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4.8), Inches(2.3))
tf = title_box.text_frame; tf.word_wrap = True
tf.text = "НАЗВАНИЕ ПРЕЗЕНТАЦИИ"
p = tf.paragraphs[0]; p.font.size = Pt(40); p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.8))
sf = sub_box.text_frame; sf.text = "Подзаголовок"
sf.paragraphs[0].font.size = Pt(20); sf.paragraphs[0].font.color.rgb = RGBColor(253, 186, 116)

# ── Контентный слайд ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(12, 10, 9)

sidebar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.22), Inches(7.5))
sidebar.fill.solid(); sidebar.fill.fore_color.rgb = RGBColor(234, 88, 12); sidebar.line.fill.background()

header = slide.shapes.add_shape(1, Inches(0.22), Inches(0), Inches(9.78), Inches(1.15))
header.fill.solid(); header.fill.fore_color.rgb = RGBColor(28, 25, 23); header.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.13), Inches(8.5), Inches(0.88))
tf = title_box.text_frame; tf.text = "ЗАГОЛОВОК СЛАЙДА"
p = tf.paragraphs[0]; p.font.size = Pt(28); p.font.bold = True
p.font.color.rgb = RGBColor(253, 186, 116)

text_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(4.8), Inches(5.9))
tf2 = text_box.text_frame; tf2.word_wrap = True
tf2.text = "Первый пункт"
tf2.paragraphs[0].font.size = Pt(16); tf2.paragraphs[0].font.color.rgb = RGBColor(231, 229, 228)
for txt in ["Второй пункт", "Третий пункт"]:
    p = tf2.add_paragraph(); p.text = txt
    p.font.size = Pt(16); p.font.color.rgb = RGBColor(231, 229, 228); p.space_before = Pt(12)

if os.path.exists("img_1.jpg"):
    slide.shapes.add_picture("img_1.jpg", Inches(5.7), Inches(1.2), width=Inches(4.1), height=Inches(6.0))

prs.save("presentation.pptx")
"""

TEMPLATE5 = """
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ── Титульный слайд (картинка на весь фон, текст поверх снизу) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(10, 14, 26)

if os.path.exists("img_0.jpg"):
    slide.shapes.add_picture("img_0.jpg", Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

# Тёмная плашка снизу поверх картинки
overlay = slide.shapes.add_shape(1, Inches(0), Inches(4.2), Inches(10), Inches(3.3))
overlay.fill.solid(); overlay.fill.fore_color.rgb = RGBColor(10, 14, 26); overlay.line.fill.background()

# Акцентная линия над плашкой
accent_line = slide.shapes.add_shape(1, Inches(0), Inches(4.18), Inches(10), Inches(0.06))
accent_line.fill.solid(); accent_line.fill.fore_color.rgb = RGBColor(56, 189, 248); accent_line.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.4), Inches(8.6), Inches(1.8))
tf = title_box.text_frame; tf.word_wrap = True
tf.text = "НАЗВАНИЕ ПРЕЗЕНТАЦИИ"
p = tf.paragraphs[0]; p.font.size = Pt(44); p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.3), Inches(8), Inches(0.7))
sf = sub_box.text_frame; sf.text = "Подзаголовок"
sf.paragraphs[0].font.size = Pt(18); sf.paragraphs[0].font.color.rgb = RGBColor(56, 189, 248)

# ── Контентный слайд (картинка справа как вертикальная панель) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(10, 14, 26)

header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.1))
header.fill.solid(); header.fill.fore_color.rgb = RGBColor(17, 24, 39); header.line.fill.background()

accent_left = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), Inches(1.1))
accent_left.fill.solid(); accent_left.fill.fore_color.rgb = RGBColor(56, 189, 248); accent_left.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.12), Inches(8.5), Inches(0.85))
tf = title_box.text_frame; tf.text = "ЗАГОЛОВОК СЛАЙДА"
p = tf.paragraphs[0]; p.font.size = Pt(27); p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.0), Inches(5.9))
tf2 = text_box.text_frame; tf2.word_wrap = True
tf2.text = "Первый пункт"
tf2.paragraphs[0].font.size = Pt(16); tf2.paragraphs[0].font.color.rgb = RGBColor(226, 232, 240)
for txt in ["Второй пункт", "Третий пункт"]:
    p = tf2.add_paragraph(); p.text = txt
    p.font.size = Pt(16); p.font.color.rgb = RGBColor(226, 232, 240); p.space_before = Pt(14)

if os.path.exists("img_1.jpg"):
    slide.shapes.add_picture("img_1.jpg", Inches(5.8), Inches(1.15), width=Inches(4.2), height=Inches(6.35))

# Градиентная плашка-переход слева от картинки
fade = slide.shapes.add_shape(1, Inches(5.6), Inches(1.15), Inches(0.4), Inches(6.35))
fade.fill.solid(); fade.fill.fore_color.rgb = RGBColor(10, 14, 26); fade.line.fill.background()

prs.save("presentation.pptx")
"""

TEMPLATE6 = """
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ── Титульный слайд (бордово-золотой, картинка центр-верх) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(28, 8, 20)

# Верхняя полоса с картинкой
top_band = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(4.2))
top_band.fill.solid(); top_band.fill.fore_color.rgb = RGBColor(40, 10, 28); top_band.line.fill.background()

if os.path.exists("img_0.jpg"):
    slide.shapes.add_picture("img_0.jpg", Inches(2.5), Inches(0.2), width=Inches(5.0), height=Inches(3.6))

# Тёмный fade снизу картинки
fade_bot = slide.shapes.add_shape(1, Inches(0), Inches(3.4), Inches(10), Inches(1.0))
fade_bot.fill.solid(); fade_bot.fill.fore_color.rgb = RGBColor(28, 8, 20); fade_bot.line.fill.background()

# Золотая линия-разделитель
gold_line = slide.shapes.add_shape(1, Inches(1.5), Inches(4.25), Inches(7.0), Inches(0.06))
gold_line.fill.solid(); gold_line.fill.fore_color.rgb = RGBColor(212, 175, 55); gold_line.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.45), Inches(8.6), Inches(1.7))
tf = title_box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "НАЗВАНИЕ ПРЕЗЕНТАЦИИ"
p.font.size = Pt(42); p.font.bold = True; p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = RGBColor(255, 255, 255)

sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.25), Inches(8.6), Inches(0.7))
sf = sub_box.text_frame
sp = sf.paragraphs[0]; sp.text = "Подзаголовок"; sp.alignment = PP_ALIGN.CENTER
sp.font.size = Pt(18); sp.font.color.rgb = RGBColor(212, 175, 55)

# ── Контентный слайд ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(28, 8, 20)

# Картинка как широкая горизонтальная полоса сверху
if os.path.exists("img_1.jpg"):
    slide.shapes.add_picture("img_1.jpg", Inches(0), Inches(0), width=Inches(10), height=Inches(2.8))

# Тёмный overlay поверх картинки (только нижняя часть полосы) для читаемости заголовка
header_overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(2.8))
header_overlay.fill.solid(); header_overlay.fill.fore_color.rgb = RGBColor(28, 8, 20); header_overlay.line.fill.background()

# Золотая линия под картинкой
gold_line2 = slide.shapes.add_shape(1, Inches(0), Inches(2.78), Inches(10), Inches(0.06))
gold_line2.fill.solid(); gold_line2.fill.fore_color.rgb = RGBColor(212, 175, 55); gold_line2.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.8), Inches(1.8))
tf = title_box.text_frame; tf.word_wrap = True
tf.text = "ЗАГОЛОВОК СЛАЙДА"
p = tf.paragraphs[0]; p.font.size = Pt(30); p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

text_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(8.8), Inches(4.2))
tf2 = text_box.text_frame; tf2.word_wrap = True
tf2.text = "Первый пункт"
tf2.paragraphs[0].font.size = Pt(17); tf2.paragraphs[0].font.color.rgb = RGBColor(245, 230, 210)
for txt in ["Второй пункт", "Третий пункт"]:
    p = tf2.add_paragraph(); p.text = txt
    p.font.size = Pt(17); p.font.color.rgb = RGBColor(245, 230, 210); p.space_before = Pt(14)

prs.save("presentation.pptx")
"""

TEMPLATES_PREMIUM = [TEMPLATE1, TEMPLATE2, TEMPLATE3, TEMPLATE4, TEMPLATE5, TEMPLATE6]

# ==================== ШАБЛОНЫ ДЛЯ БЕСПЛАТНЫХ (без картинок) ====================
TEMPLATE_FREE_1 = """
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(17, 24, 39)

panel = slide.shapes.add_shape(1, Inches(0.6), Inches(1.8), Inches(8.8), Inches(3.4))
panel.fill.solid(); panel.fill.fore_color.rgb = RGBColor(31, 41, 55)
panel.line.color.rgb = RGBColor(99, 102, 241); panel.line.width = Pt(1.5)

title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(8.0), Inches(1.6))
tf = title_box.text_frame
tf.text = "НАЗВАНИЕ ПРЕЗЕНТАЦИИ"
p = tf.paragraphs[0]; p.font.size = Pt(42); p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255); p.alignment = PP_ALIGN.CENTER

sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(8.0), Inches(0.7))
sf = sub_box.text_frame; sf.text = "Подзаголовок"
sf.paragraphs[0].font.size = Pt(18); sf.paragraphs[0].font.color.rgb = RGBColor(165, 180, 252)
sf.paragraphs[0].alignment = PP_ALIGN.CENTER

slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(17, 24, 39)

header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.3))
header.fill.solid(); header.fill.fore_color.rgb = RGBColor(99, 102, 241); header.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.9))
tf = title_box.text_frame; tf.text = "ЗАГОЛОВОК СЛАЙДА"
p = tf.paragraphs[0]; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255)

text_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.8), Inches(5.6))
tf2 = text_box.text_frame; tf2.word_wrap = True
tf2.text = "Первый пункт"
tf2.paragraphs[0].font.size = Pt(17); tf2.paragraphs[0].font.color.rgb = RGBColor(229, 231, 235)
for txt in ["Второй пункт", "Третий пункт"]:
    p = tf2.add_paragraph(); p.text = txt
    p.font.size = Pt(17); p.font.color.rgb = RGBColor(229, 231, 235); p.space_before = Pt(12)

prs.save("presentation.pptx")
"""

TEMPLATE_FREE_2 = """
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(255, 255, 255)

left_panel = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(4.2), Inches(7.5))
left_panel.fill.solid(); left_panel.fill.fore_color.rgb = RGBColor(5, 150, 105); left_panel.line.fill.background()

deco = slide.shapes.add_shape(1, Inches(0), Inches(5.5), Inches(4.2), Inches(2.0))
deco.fill.solid(); deco.fill.fore_color.rgb = RGBColor(4, 120, 87); deco.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.8), Inches(3.6), Inches(2.5))
tf = title_box.text_frame; tf.word_wrap = True
tf.text = "НАЗВАНИЕ ПРЕЗЕНТАЦИИ"
p = tf.paragraphs[0]; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255)

sub_box = slide.shapes.add_textbox(Inches(0.3), Inches(4.5), Inches(3.6), Inches(0.8))
sf = sub_box.text_frame; sf.text = "Подзаголовок"
sf.paragraphs[0].font.size = Pt(17); sf.paragraphs[0].font.color.rgb = RGBColor(167, 243, 208)

slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(255, 255, 255)

left_panel = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.2), Inches(7.5))
left_panel.fill.solid(); left_panel.fill.fore_color.rgb = RGBColor(5, 150, 105); left_panel.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.85))
tf = title_box.text_frame; tf.text = "ЗАГОЛОВОК СЛАЙДА"
p = tf.paragraphs[0]; p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = RGBColor(5,150,105)

sep = slide.shapes.add_shape(1, Inches(0.5), Inches(1.25), Inches(3), Inches(0.07))
sep.fill.solid(); sep.fill.fore_color.rgb = RGBColor(5,150,105); sep.line.fill.background()

text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.7))
tf2 = text_box.text_frame; tf2.word_wrap = True
tf2.text = "Первый пункт"
tf2.paragraphs[0].font.size = Pt(17); tf2.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)
for txt in ["Второй пункт", "Третий пункт"]:
    p = tf2.add_paragraph(); p.text = txt
    p.font.size = Pt(17); p.font.color.rgb = RGBColor(31, 41, 55); p.space_before = Pt(14)

prs.save("presentation.pptx")
"""

TEMPLATE_FREE_3 = """
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(255, 247, 237)

top_block = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(4.2))
top_block.fill.solid(); top_block.fill.fore_color.rgb = RGBColor(234, 88, 12); top_block.line.fill.background()

accent_block = slide.shapes.add_shape(1, Inches(0), Inches(3.8), Inches(10), Inches(0.6))
accent_block.fill.solid(); accent_block.fill.fore_color.rgb = RGBColor(194, 65, 12); accent_block.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.1), Inches(8.6), Inches(1.8))
tf = title_box.text_frame
tf.text = "НАЗВАНИЕ ПРЕЗЕНТАЦИИ"
p = tf.paragraphs[0]; p.font.size = Pt(44); p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255); p.alignment = PP_ALIGN.CENTER

sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(8.6), Inches(0.7))
sf = sub_box.text_frame; sf.text = "Подзаголовок"
sf.paragraphs[0].font.size = Pt(19); sf.paragraphs[0].font.color.rgb = RGBColor(254, 215, 170)
sf.paragraphs[0].alignment = PP_ALIGN.CENTER

slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(255, 247, 237)

header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
header.fill.solid(); header.fill.fore_color.rgb = RGBColor(234, 88, 12); header.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.85))
tf = title_box.text_frame; tf.text = "ЗАГОЛОВОК СЛАЙДА"
p = tf.paragraphs[0]; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255)

for i, (x, txt, color) in enumerate([
    (Inches(0.4), "Первый пункт", RGBColor(234,88,12)),
    (Inches(3.55), "Второй пункт", RGBColor(194,65,12)),
    (Inches(6.7), "Третий пункт", RGBColor(154,52,18)),
]):
    card = slide.shapes.add_shape(1, x, Inches(1.5), Inches(2.9), Inches(5.5))
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(255,255,255)
    card.line.color.rgb = color; card.line.width = Pt(1.5)
    top = slide.shapes.add_shape(1, x, Inches(1.5), Inches(2.9), Inches(0.18))
    top.fill.solid(); top.fill.fore_color.rgb = color; top.line.fill.background()
    tb = slide.shapes.add_textbox(x + Inches(0.15), Inches(1.85), Inches(2.6), Inches(4.6))
    tf2 = tb.text_frame; tf2.word_wrap = True
    tf2.text = txt; tf2.paragraphs[0].font.size = Pt(15)
    tf2.paragraphs[0].font.color.rgb = RGBColor(31,41,55); tf2.paragraphs[0].font.bold = True

prs.save("presentation.pptx")
"""

TEMPLATE_FREE_4 = """
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(10, 10, 20)

for i in range(5):
    line = slide.shapes.add_shape(1, Inches(0), Inches(i * 1.5), Inches(10), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(80, 40, 120); line.line.fill.background()

glow = slide.shapes.add_shape(1, Inches(1.5), Inches(1.5), Inches(7), Inches(4))
glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(30, 10, 60)
glow.line.color.rgb = RGBColor(167, 139, 250); glow.line.width = Pt(1)

title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(8.0), Inches(1.8))
tf = title_box.text_frame
tf.text = "НАЗВАНИЕ ПРЕЗЕНТАЦИИ"
p = tf.paragraphs[0]; p.font.size = Pt(42); p.font.bold = True
p.font.color.rgb = RGBColor(216, 180, 254); p.alignment = PP_ALIGN.CENTER

sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.9), Inches(8.0), Inches(0.7))
sf = sub_box.text_frame; sf.text = "Подзаголовок"
sf.paragraphs[0].font.size = Pt(18); sf.paragraphs[0].font.color.rgb = RGBColor(196, 181, 253)
sf.paragraphs[0].alignment = PP_ALIGN.CENTER

slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(10, 10, 20)

for i in range(5):
    line = slide.shapes.add_shape(1, Inches(0), Inches(i * 1.5), Inches(10), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(80,40,120); line.line.fill.background()

header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
header.fill.solid(); header.fill.fore_color.rgb = RGBColor(109,40,217); header.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.85))
tf = title_box.text_frame; tf.text = "ЗАГОЛОВОК СЛАЙДА"
p = tf.paragraphs[0]; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255)

text_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.8), Inches(5.7))
tf2 = text_box.text_frame; tf2.word_wrap = True
tf2.text = "Первый пункт"
tf2.paragraphs[0].font.size = Pt(17); tf2.paragraphs[0].font.color.rgb = RGBColor(216, 180, 254)
for txt in ["Второй пункт", "Третий пункт"]:
    p = tf2.add_paragraph(); p.text = txt
    p.font.size = Pt(17); p.font.color.rgb = RGBColor(216, 180, 254); p.space_before = Pt(14)

prs.save("presentation.pptx")
"""

TEMPLATE_FREE_5 = """
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(248, 250, 252)

side = slide.shapes.add_shape(1, Inches(6.8), Inches(0), Inches(3.2), Inches(7.5))
side.fill.solid(); side.fill.fore_color.rgb = RGBColor(30, 64, 175); side.line.fill.background()

tri_top = slide.shapes.add_shape(1, Inches(5.8), Inches(0), Inches(1.2), Inches(7.5))
tri_top.fill.solid(); tri_top.fill.fore_color.rgb = RGBColor(59, 130, 246); tri_top.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(5.8), Inches(2.0))
tf = title_box.text_frame; tf.word_wrap = True
tf.text = "НАЗВАНИЕ ПРЕЗЕНТАЦИИ"
p = tf.paragraphs[0]; p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = RGBColor(15,23,42)

sep = slide.shapes.add_shape(1, Inches(0.5), Inches(4.0), Inches(2.5), Inches(0.07))
sep.fill.solid(); sep.fill.fore_color.rgb = RGBColor(59,130,246); sep.line.fill.background()

sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(5.8), Inches(0.7))
sf = sub_box.text_frame; sf.text = "Подзаголовок"
sf.paragraphs[0].font.size = Pt(18); sf.paragraphs[0].font.color.rgb = RGBColor(100, 116, 139)

slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(248, 250, 252)

side = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.15))
side.fill.solid(); side.fill.fore_color.rgb = RGBColor(30, 64, 175); side.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.9))
tf = title_box.text_frame; tf.text = "ЗАГОЛОВОК СЛАЙДА"
p = tf.paragraphs[0]; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255)

text_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(8.8), Inches(5.8))
tf2 = text_box.text_frame; tf2.word_wrap = True
tf2.text = "Первый пункт"
tf2.paragraphs[0].font.size = Pt(17); tf2.paragraphs[0].font.color.rgb = RGBColor(15, 23, 42)
for txt in ["Второй пункт", "Третий пункт"]:
    p = tf2.add_paragraph(); p.text = txt
    p.font.size = Pt(17); p.font.color.rgb = RGBColor(15, 23, 42); p.space_before = Pt(14)

prs.save("presentation.pptx")
"""

TEMPLATES_FREE = [TEMPLATE_FREE_1, TEMPLATE_FREE_2, TEMPLATE_FREE_3, TEMPLATE_FREE_4, TEMPLATE_FREE_5]
